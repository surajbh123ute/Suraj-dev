import os
import subprocess
from pptx import Presentation
from llama_index.core import Document
from document_loaders.base_loader import SimpleDocumentLoader
from core.enums import DocumentType
from utils import process_graph, is_graph

class PPTXDocumentLoader(SimpleDocumentLoader):
    """Loader for PPT and PPTX documents."""
    def _load_data(self, ppt_file) -> list[Document]:
        ppt_path = ppt_file.name  # Assuming 'file' object has a 'name' attribute
        pdf_path = self._convert_ppt_to_pdf(ppt_path)
        images_data = self._convert_pdf_to_images(pdf_path)
        slide_texts = self._extract_text_and_notes_from_ppt(ppt_path)
        processed_data = []

        for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
            if notes:
                notes = "\n\nThe speaker notes for this slide are: " + notes

            with open(image_path, 'rb') as image_file:
                image_content = image_file.read()

            image_description = " "
            if is_graph(image_content):
                image_description = process_graph(image_content)

            image_metadata = {
                "source": f"{os.path.basename(ppt_path)}",
                "image": image_path,
                "caption": slide_text + image_description + notes,
                "type": DocumentType.IMAGE.value,
                "page_num": page_num
            }
            processed_data.append(Document(text="This is a slide with the text: " + slide_text + image_description, metadata=image_metadata))

        return processed_data

    def _convert_ppt_to_pdf(self, ppt_path):
        """Convert a PowerPoint file to PDF using LibreOffice."""
        base_name = os.path.basename(ppt_path)
        ppt_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
        new_dir_path = os.path.abspath("vectorstore/ppt_references")
        os.makedirs(new_dir_path, exist_ok=True)
        pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")
        command = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', new_dir_path, ppt_path]
        subprocess.run(command, check=True)
        return pdf_path

    def _convert_pdf_to_images(self, pdf_path):
        """Convert a PDF file to a series of images using PyMuPDF."""
        import fitz
        doc = fitz.open(pdf_path)
        base_name = os.path.basename(pdf_path)
        pdf_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
        new_dir_path = os.path.join(os.getcwd(), "vectorstore/ppt_references")
        os.makedirs(new_dir_path, exist_ok=True)
        image_paths = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            output_image_path = os.path.join(new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png")
            pix.save(output_image_path)
            image_paths.append((output_image_path, page_num))
        doc.close()
        return image_paths

    def _extract_text_and_notes_from_ppt(self, ppt_path):
        """Extract text and notes from a PowerPoint file."""
        prs = Presentation(ppt_path)
        text_and_notes = []
        for slide in prs.slides:
            slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            try:
                notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
            except:
                notes = ''
            text_and_notes.append((slide_text, notes))
        return text_and_notes